"""
artifacts_service.py — TraceLearn.ai
Generates PDF report, PPTX presentation, and summary PDF from AI analysis data.
Uploads all three to S3 and returns an ArtifactsResponse.
"""

from __future__ import annotations

import io
import os
import textwrap
from datetime import datetime

import boto3
from botocore.exceptions import BotoCoreError, ClientError

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    HRFlowable, ListFlowable, ListItem, PageBreak,
    Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from models.ai_models import ArtifactsRequest, ArtifactsResponse

# ─── Colour palette ───────────────────────────────────────────────────────────

C_PRIMARY     = colors.HexColor("#4F46E5")
C_PRIMARY_LT  = colors.HexColor("#EEF2FF")
C_DESTRUCTIVE = colors.HexColor("#DC2626")
C_SUCCESS     = colors.HexColor("#059669")
C_DARK        = colors.HexColor("#111827")
C_MID         = colors.HexColor("#6B7280")
C_LIGHT       = colors.HexColor("#F9FAFB")
C_BORDER      = colors.HexColor("#E5E7EB")
C_CODE_BG     = colors.HexColor("#1E1E2E")
C_CODE_FG     = colors.HexColor("#CDD6F4")

PPT_BG      = RGBColor(0x0F, 0x0F, 0x23)
PPT_PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
PPT_ACCENT  = RGBColor(0x06, 0xB6, 0xD4)
PPT_TEXT    = RGBColor(0xF9, 0xFA, 0xFB)
PPT_MUTED   = RGBColor(0x9C, 0xA3, 0xAF)
PPT_SUCCESS = RGBColor(0x05, 0x96, 0x69)
PPT_ERROR   = RGBColor(0xDC, 0x26, 0x26)


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _safe(val, fallback: str = "—") -> str:
    if val is None:
        return fallback
    s = str(val).strip()
    return s if s and s not in ("None", "null", "undefined") else fallback

def _now() -> str:
    return datetime.utcnow().strftime("%B %d, %Y · %H:%M UTC")


# ─── PDF shared styles ────────────────────────────────────────────────────────

def _pdf_styles():
    def s(name, **kw):
        return ParagraphStyle(name, **kw)
    return {
        "cover_title": s("cover_title", fontName="Helvetica-Bold", fontSize=28,
            textColor=C_PRIMARY, leading=34, alignment=TA_CENTER, spaceAfter=8),
        "cover_sub": s("cover_sub", fontName="Helvetica", fontSize=13,
            textColor=C_MID, leading=18, alignment=TA_CENTER, spaceAfter=4),
        "cover_meta": s("cover_meta", fontName="Helvetica", fontSize=10,
            textColor=C_MID, alignment=TA_CENTER),
        "h1": s("h1", fontName="Helvetica-Bold", fontSize=18, textColor=C_DARK,
            spaceBefore=18, spaceAfter=8, leading=24),
        "h2": s("h2", fontName="Helvetica-Bold", fontSize=13, textColor=C_PRIMARY,
            spaceBefore=14, spaceAfter=6, leading=18),
        "body": s("body", fontName="Helvetica", fontSize=10, textColor=C_DARK,
            leading=16, spaceAfter=6),
        "body_muted": s("body_muted", fontName="Helvetica", fontSize=9,
            textColor=C_MID, leading=14, spaceAfter=4),
        "code": s("code", fontName="Courier", fontSize=8.5, textColor=C_CODE_FG,
            backColor=C_CODE_BG, leading=13, spaceAfter=4,
            leftIndent=8, rightIndent=8, borderPad=6),
        "bullet": s("bullet", fontName="Helvetica", fontSize=10, textColor=C_DARK,
            leading=16, leftIndent=14, spaceAfter=3),
        "label": s("label", fontName="Helvetica-Bold", fontSize=9, textColor=C_PRIMARY,
            leading=14, spaceAfter=2),
        "footer": s("footer", fontName="Helvetica", fontSize=8,
            textColor=C_MID, alignment=TA_CENTER),
    }

def _divider():
    return HRFlowable(width="100%", thickness=1, color=C_BORDER, spaceAfter=8, spaceBefore=4)

def _section_heading(text, st):
    return [Spacer(1, 0.15 * cm), Paragraph(text, st["h1"]), _divider()]

def _info_table(rows, st):
    data = [[Paragraph(k, st["label"]), Paragraph(v, st["body"])] for k, v in rows]
    t = Table(data, colWidths=[4.5 * cm, 12 * cm])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, -1), C_PRIMARY_LT),
        ("TEXTCOLOR",  (0, 0), (0, -1), C_PRIMARY),
        ("ALIGN",      (0, 0), (-1, -1), "LEFT"),
        ("VALIGN",     (0, 0), (-1, -1), "TOP"),
        ("FONTNAME",   (0, 0), (0, -1), "Helvetica-Bold"),
        ("FONTSIZE",   (0, 0), (-1, -1), 9),
        ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, C_LIGHT]),
        ("GRID",       (0, 0), (-1, -1), 0.5, C_BORDER),
        ("LEFTPADDING",  (0, 0), (-1, -1), 8),
        ("RIGHTPADDING", (0, 0), (-1, -1), 8),
        ("TOPPADDING",   (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING",(0, 0), (-1, -1), 6),
    ]))
    return t

def _code_block(code, st):
    escaped = _safe(code, "(no code)")[:3000].replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
    return [Paragraph(escaped, st["code"])]


# ─── 1. Error Explanation Report (PDF) ───────────────────────────────────────

def generate_error_report_pdf(p: ArtifactsRequest) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm,
        title="Error Explanation Report", author="TraceLearn.ai")
    st = _pdf_styles()
    story = []

    story += [
        Spacer(1, 1.5*cm),
        Paragraph("TraceLearn.ai", st["cover_sub"]),
        Spacer(1, 0.4*cm),
        Paragraph("Error Explanation Report", st["cover_title"]),
        Spacer(1, 0.3*cm),
        Paragraph(_safe(p.explanation, "Code Error Analysis"), st["cover_sub"]),
        Spacer(1, 0.5*cm),
        Paragraph(f"Session: {_safe(p.sessionId)}  ·  Generated: {_now()}", st["cover_meta"]),
        Spacer(1, 1.5*cm),
        HRFlowable(width="100%", thickness=3, color=C_PRIMARY, spaceAfter=12),
    ]

    story += _section_heading("1 · Error Summary", st)
    story += [_info_table([
        ("Error Type", _safe(p.explanation, "Unknown").split(":")[0][:60]),
        ("Session ID",  _safe(p.sessionId)),
        ("Generated",   _now()),
    ], st), Spacer(1, 0.4*cm)]

    story += _section_heading("2 · Root Cause", st)
    story += [Paragraph(_safe(p.whyItHappened, "Not available."), st["body"]), Spacer(1, 0.3*cm)]

    if p.conceptBehindError:
        story += [Paragraph("Core Concept", st["h2"]),
                  Paragraph(_safe(p.conceptBehindError), st["body"]), Spacer(1, 0.3*cm)]

    if p.stepByStepReasoning:
        story += _section_heading("3 · Step-by-Step Reasoning", st)
        items = [ListItem(Paragraph(f"<b>Step {i+1}:</b> {_safe(step)}", st["bullet"]),
                          leftIndent=16)
                 for i, step in enumerate(p.stepByStepReasoning)]
        story += [ListFlowable(items, bulletType="bullet", start="•"), Spacer(1, 0.3*cm)]

    if p.code:
        story += _section_heading("4 · Original Code", st)
        story += _code_block(p.code, st) + [Spacer(1, 0.3*cm)]

    if p.fixedCode:
        story += _section_heading("5 · Fixed Code", st)
        story += _code_block(p.fixedCode, st) + [Spacer(1, 0.3*cm)]

    if p.fixAnalysis:
        story += _section_heading("6 · Fix Analysis", st)
        rows = []
        if p.fixAnalysis.whatChanged:    rows.append(("What Changed",       _safe(p.fixAnalysis.whatChanged)))
        if p.fixAnalysis.whyItWorks:     rows.append(("Why It Works",       _safe(p.fixAnalysis.whyItWorks)))
        if p.fixAnalysis.reinforcedConcept: rows.append(("Reinforced Concept", _safe(p.fixAnalysis.reinforcedConcept)))
        if rows:
            story += [_info_table(rows, st), Spacer(1, 0.3*cm)]

    if p.learningResources:
        story += _section_heading("7 · Learning Resources", st)
        for r in p.learningResources:
            story += [
                Paragraph(f"<b>{_safe(r.title)}</b>", st["body"]),
                Paragraph(f'<link href="{_safe(r.url)}" color="#4F46E5">{_safe(r.url)}</link>', st["body_muted"]),
                Spacer(1, 0.2*cm),
            ]

    story += [
        PageBreak(), Spacer(1, 2*cm),
        HRFlowable(width="100%", thickness=1, color=C_BORDER), Spacer(1, 0.2*cm),
        Paragraph("Generated by TraceLearn.ai — AI-powered developer learning platform", st["footer"]),
    ]

    doc.build(story)
    return buf.getvalue()


# ─── 2. Learning Presentation (PPTX) ─────────────────────────────────────────

def _pptx_bg(slide, prs):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PPT_BG

def _pptx_add_text(tf, text, size, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or PPT_TEXT

def _accent_bar(slide, prs):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPT_PRIMARY
    bar.line.fill.background()

def generate_presentation_pptx(p: ArtifactsRequest) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide():
        sl = prs.slides.add_slide(blank)
        _pptx_bg(sl, prs)
        _accent_bar(sl, prs)
        return sl

    def add_box(slide, l, t, w, h, bg=None):
        box = slide.shapes.add_shape(1, l, t, w, h)
        if bg: box.fill.solid(); box.fill.fore_color.rgb = bg
        else:  box.fill.background()
        box.line.fill.background()
        return box

    def tb(slide, l, t, w, h):
        return slide.shapes.add_textbox(l, t, w, h)

    # Slide 1: Title
    sl = new_slide()
    c = sl.shapes.add_shape(9, Inches(10.5), Inches(-1), Inches(4), Inches(4))
    c.fill.solid(); c.fill.fore_color.rgb = PPT_PRIMARY; c.line.fill.background()
    tf = tb(sl, Inches(0.7), Inches(1.8), Inches(9), Inches(1.2)).text_frame
    tf.word_wrap = True
    _pptx_add_text(tf, "Error Explanation", 40, bold=True)
    tf2 = tb(sl, Inches(0.7), Inches(2.9), Inches(9), Inches(0.7)).text_frame
    tf2.word_wrap = True
    _pptx_add_text(tf2, _safe(p.explanation, "Code Error Analysis")[:80], 18, color=PPT_MUTED)
    tf3 = tb(sl, Inches(0.7), Inches(6.95), Inches(8), Inches(0.5)).text_frame
    _pptx_add_text(tf3, f"TraceLearn.ai  ·  Session {_safe(p.sessionId)[:20]}  ·  {_now()}", 10, color=PPT_MUTED)

    # Slide 2: What Went Wrong
    sl = new_slide()
    add_box(sl, Inches(0), Inches(0.08), Inches(0.25), Inches(7.42), bg=PPT_ERROR)
    _pptx_add_text(tb(sl, Inches(0.7), Inches(0.3), Inches(11), Inches(0.7)).text_frame, "What Went Wrong", 28, bold=True, color=PPT_ERROR)
    tf = tb(sl, Inches(0.7), Inches(1.2), Inches(11.5), Inches(5)).text_frame
    tf.word_wrap = True
    _pptx_add_text(tf, _safe(p.explanation, "The code produced an error.")[:600], 16)

    # Slide 3: Root Cause
    sl = new_slide()
    add_box(sl, Inches(0), Inches(0.08), Inches(0.25), Inches(7.42), bg=PPT_PRIMARY)
    _pptx_add_text(tb(sl, Inches(0.7), Inches(0.3), Inches(11), Inches(0.7)).text_frame, "Root Cause", 28, bold=True, color=PPT_PRIMARY)
    tf = tb(sl, Inches(0.7), Inches(1.2), Inches(11.5), Inches(3)).text_frame
    tf.word_wrap = True
    _pptx_add_text(tf, _safe(p.whyItHappened, "Root cause not available.")[:500], 15)
    if p.conceptBehindError:
        add_box(sl, Inches(0.7), Inches(4.5), Inches(11.5), Inches(1.5), bg=RGBColor(0x1E, 0x1E, 0x3F))
        _pptx_add_text(tb(sl, Inches(1.0), Inches(4.6), Inches(11), Inches(0.5)).text_frame, "Core Concept", 11, bold=True, color=PPT_ACCENT)
        tf2 = tb(sl, Inches(1.0), Inches(5.1), Inches(11), Inches(0.8)).text_frame
        tf2.word_wrap = True
        _pptx_add_text(tf2, _safe(p.conceptBehindError)[:120], 14)

    # Slide 4: Step-by-Step
    if p.stepByStepReasoning:
        sl = new_slide()
        add_box(sl, Inches(0), Inches(0.08), Inches(0.25), Inches(7.42), bg=PPT_ACCENT)
        _pptx_add_text(tb(sl, Inches(0.7), Inches(0.3), Inches(11), Inches(0.7)).text_frame, "Step-by-Step Reasoning", 28, bold=True, color=PPT_ACCENT)
        top = Inches(1.3)
        for i, step in enumerate(p.stepByStepReasoning[:6]):
            add_box(sl, Inches(0.7), top, Inches(0.5), Inches(0.5), bg=PPT_PRIMARY)
            _pptx_add_text(tb(sl, Inches(0.7), top+Inches(0.05), Inches(0.5), Inches(0.45)).text_frame, str(i+1), 13, bold=True, align=PP_ALIGN.CENTER)
            stf = tb(sl, Inches(1.4), top, Inches(11), Inches(0.55)).text_frame
            stf.word_wrap = True
            _pptx_add_text(stf, _safe(step)[:120], 13)
            top += Inches(0.8)

    # Slide 5: The Fix
    sl = new_slide()
    add_box(sl, Inches(0), Inches(0.08), Inches(0.25), Inches(7.42), bg=PPT_SUCCESS)
    _pptx_add_text(tb(sl, Inches(0.7), Inches(0.3), Inches(11), Inches(0.7)).text_frame, "The Fix", 28, bold=True, color=PPT_SUCCESS)
    if p.fixAnalysis:
        top = Inches(1.2)
        for label, value in [
            ("What Changed", _safe(p.fixAnalysis.whatChanged)),
            ("Why It Works", _safe(p.fixAnalysis.whyItWorks)),
            ("Reinforced Concept", _safe(p.fixAnalysis.reinforcedConcept)),
        ]:
            _pptx_add_text(tb(sl, Inches(0.7), top, Inches(3.5), Inches(0.4)).text_frame, label, 11, bold=True, color=PPT_ACCENT)
            vtf = tb(sl, Inches(4.3), top, Inches(8.5), Inches(0.6)).text_frame
            vtf.word_wrap = True
            _pptx_add_text(vtf, value[:180], 12)
            top += Inches(1.0)
    elif p.fixedCode:
        tf = tb(sl, Inches(0.7), Inches(1.2), Inches(11.5), Inches(5)).text_frame
        tf.word_wrap = True
        _pptx_add_text(tf, _safe(p.fixedCode)[:400], 11, color=RGBColor(0xCD, 0xD6, 0xF4))

    # Slide 6: Learning Resources
    if p.learningResources:
        sl = new_slide()
        add_box(sl, Inches(0), Inches(0.08), Inches(0.25), Inches(7.42), bg=RGBColor(0xD9, 0x77, 0x06))
        _pptx_add_text(tb(sl, Inches(0.7), Inches(0.3), Inches(11), Inches(0.7)).text_frame, "Learning Resources", 28, bold=True, color=RGBColor(0xD9, 0x77, 0x06))
        top = Inches(1.3)
        for r in p.learningResources[:5]:
            add_box(sl, Inches(0.7), top, Inches(11.5), Inches(0.7), bg=RGBColor(0x1E, 0x1E, 0x3F))
            _pptx_add_text(tb(sl, Inches(1.0), top+Inches(0.05), Inches(6), Inches(0.35)).text_frame, _safe(r.title)[:60], 13, bold=True)
            _pptx_add_text(tb(sl, Inches(1.0), top+Inches(0.38), Inches(10), Inches(0.28)).text_frame, _safe(r.url)[:80], 10, color=PPT_ACCENT)
            top += Inches(0.9)

    # Slide 7: Key Takeaway
    sl = new_slide()
    add_box(sl, Inches(0), Inches(0.08), prs.slide_width, Inches(7.42), bg=RGBColor(0x1E, 0x1E, 0x3F))
    _accent_bar(sl, prs)
    _pptx_add_text(tb(sl, Inches(1), Inches(1.5), Inches(11), Inches(0.8)).text_frame, "Key Takeaway", 34, bold=True, color=PPT_PRIMARY, align=PP_ALIGN.CENTER)
    tf = tb(sl, Inches(1), Inches(2.6), Inches(11), Inches(2.5)).text_frame
    tf.word_wrap = True
    _pptx_add_text(tf, _safe(p.learningSummary, "Review your error analysis for key takeaways.")[:300], 18, align=PP_ALIGN.CENTER)
    _pptx_add_text(tb(sl, Inches(1), Inches(6.4), Inches(11), Inches(0.5)).text_frame, "TraceLearn.ai  ·  AI-powered developer learning", 11, color=PPT_MUTED, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── 3. Session Learning Summary (PDF) ───────────────────────────────────────

def generate_summary_pdf(p: ArtifactsRequest) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm, topMargin=1.5*cm, bottomMargin=1.5*cm,
        title="Session Learning Summary", author="TraceLearn.ai")
    st = _pdf_styles()
    story = []

    ht = Table([[
        Paragraph("<b>TraceLearn.ai</b>", st["cover_sub"]),
        Paragraph("Session Learning Summary", st["cover_title"]),
        Paragraph(_now(), st["cover_meta"]),
    ]], colWidths=[4*cm, 9*cm, 4*cm])
    ht.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), C_PRIMARY),
        ("TEXTCOLOR",  (0, 0), (-1, -1), colors.white),
        ("ALIGN",      (0, 0), (0, 0), "LEFT"),
        ("ALIGN",      (1, 0), (1, 0), "CENTER"),
        ("ALIGN",      (2, 0), (2, 0), "RIGHT"),
        ("VALIGN",     (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING",  (0, 0), (-1, -1), 10),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ("TOPPADDING",   (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING",(0, 0), (-1, -1), 10),
        ("FONTNAME",   (0, 0), (-1, -1), "Helvetica-Bold"),
    ]))
    story += [ht, Spacer(1, 0.5*cm)]

    story += [Paragraph("Error Snapshot", st["h2"])]
    story += [_info_table([
        ("Explanation",  _safe(p.explanation)[:120]),
        ("Root Cause",   _safe(p.whyItHappened)[:120]),
        ("Core Concept", _safe(p.conceptBehindError)[:100]),
    ], st), Spacer(1, 0.4*cm)]

    story += [
        Paragraph("What You Learned", st["h2"]),
        Paragraph(_safe(p.learningSummary, "Review the error analysis for learning outcomes."), st["body"]),
        Spacer(1, 0.3*cm),
    ]

    if p.stepByStepReasoning:
        story += [Paragraph("Key Reasoning Steps", st["h2"])]
        for i, step in enumerate(p.stepByStepReasoning[:3]):
            story += [Paragraph(f"<b>{i+1}.</b> {_safe(step)}", st["bullet"])]
        story += [Spacer(1, 0.3*cm)]

    if p.fixAnalysis:
        story += [Paragraph("The Fix", st["h2"]),
                  _info_table([
                      ("What Changed",       _safe(p.fixAnalysis.whatChanged)[:120]),
                      ("Why It Works",       _safe(p.fixAnalysis.whyItWorks)[:120]),
                      ("Reinforced Concept", _safe(p.fixAnalysis.reinforcedConcept)[:100]),
                  ], st), Spacer(1, 0.3*cm)]

    if p.learningResources:
        story += [Paragraph("Further Reading", st["h2"])]
        for r in p.learningResources[:4]:
            story += [Paragraph(
                f'• <b>{_safe(r.title)}</b> — <link href="{_safe(r.url)}" color="#4F46E5">{_safe(r.url)}</link>',
                st["bullet"]
            )]
        story += [Spacer(1, 0.3*cm)]

    story += [
        _divider(),
        Paragraph(f"Session {_safe(p.sessionId)}  ·  Generated by TraceLearn.ai  ·  {_now()}", st["footer"]),
    ]

    doc.build(story)
    return buf.getvalue()


# ─── S3 Upload ────────────────────────────────────────────────────────────────

def _upload_to_s3(data: bytes, key: str, content_type: str) -> str:
    bucket = os.getenv("S3_BUCKET_NAME", "tracelearn-artifacts")
    region = os.getenv("AWS_REGION", "ap-south-1")
    client = boto3.client("s3", region_name=region)
    client.put_object(
        Bucket=bucket, Key=key, Body=data,
        ContentType=content_type, ContentDisposition="inline",
    )
    return f"https://{bucket}.s3.{region}.amazonaws.com/{key}"


# ─── Public entry point ───────────────────────────────────────────────────────

async def generate_and_upload_artifacts(payload: ArtifactsRequest) -> ArtifactsResponse:
    """
    Generates PDF report, PPTX presentation, and summary PDF.
    Uploads all three to S3.
    Returns ArtifactsResponse (NOT a dict — required by ai_service.py).
    """
    base_key = f"artifacts/{payload.sessionId}"

    pdf_bytes     = generate_error_report_pdf(payload)
    pptx_bytes    = generate_presentation_pptx(payload)
    summary_bytes = generate_summary_pdf(payload)

    try:
        pdf_url = _upload_to_s3(pdf_bytes, f"{base_key}/report.pdf", "application/pdf")
        pptx_url = _upload_to_s3(
            pptx_bytes, f"{base_key}/presentation.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        summary_url = _upload_to_s3(summary_bytes, f"{base_key}/summary.pdf", "application/pdf")
    except (BotoCoreError, ClientError) as exc:
        raise RuntimeError(f"S3 upload failed: {exc}") from exc

    # ✅ Returns ArtifactsResponse object — NOT a dict
    # ai_service.py does: return await generate_and_upload_artifacts(payload)
    # FastAPI expects ArtifactsResponse, so this must be the model, not a dict.
    return ArtifactsResponse(
        pdfUrl=pdf_url,
        pptUrl=pptx_url,
        summaryUrl=summary_url,
    )